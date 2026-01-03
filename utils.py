import shelve
import unicodedata
import os
import re

# Try importing FPDF
try:
    from fpdf import FPDF
except ImportError:
    FPDF = None

# Try importing python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except ImportError:
    Presentation = None

def clean_text(text):
    """
    Normalizes text to handle unicode characters incompatible with latin-1/ascii
    """
    return unicodedata.normalize('NFKD', text).encode('ascii', 'ignore').decode('ascii')

def generate_pdf(content, filename):
    """
    Generates a PDF file from text content.
    """
    if FPDF is None:
        return None

    try:
        content = clean_text(content)
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font('Arial', 'B', 16)
        pdf.cell(0, 10, 'Course Content', ln=True, align='C')
        pdf.ln(10)
        
        pdf.set_font('Arial', '', 12)
        pdf.multi_cell(0, 10, content)
        
        # FPDF output to file
        pdf.output(filename, 'F')
        return pdf
    except Exception as e:
        print(f"Error generating PDF: {e}")
        return None

def generate_ppt(content, filename):
    """
    Generates a PowerPoint file from markdown-like text content.
    """
    if Presentation is None:
        return None

    try:
        prs = Presentation()
        
        # Clean content
        content = clean_text(content)
        
        # Split content by Headers (assuming lines starting with # are headers)
        # matches newlines followed by #, capturing the # to keep it or just splitting
        sections = re.split(r'\n(?=#)', content)
        
        # Title Slide (First section is usually the title or intro)
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Generated Course"
        subtitle.text = "Automated Course Content Generator"

        for section in sections:
            section = section.strip()
            if not section:
                continue
                
            lines = section.split('\n')
            header = lines[0].replace('#', '').strip()
            body_text = '\n'.join(lines[1:]).strip()
            
            # Content Slide
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)
            
            shapes = slide.shapes
            title_shape = shapes.title
            body_shape = shapes.placeholders[1]
            
            title_shape.text = header[:100] # Limit title length
            body_shape.text = body_text[:1000] # Limit body length to prevent overflow (basic handling)
            
        prs.save(filename)
        return filename
    except Exception as e:
        print(f"Error generating PPT: {e}")
        return None

def load_chat_history():
    """
    Loads chat history safely.
    """
    try:
        with shelve.open("chat_history") as db:
            return db.get("messages", [])
    except Exception:
        return []

def save_chat_history(messages):
    """
    Saves chat history safely.
    """
    try:
        with shelve.open("chat_history") as db:
            db["messages"] = messages
    except Exception:
        pass